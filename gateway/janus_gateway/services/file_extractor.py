"""File extraction utilities for document attachments."""

import base64
import io
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation


class FileExtractor:
    """Extract text content from various file formats."""

    MAX_TEXT_LENGTH = 100_000

    def extract(self, content: str, mime_type: str, filename: str) -> str:
        """Extract text from file content."""
        extension = Path(filename).suffix.lower()

        if self._is_text_like(mime_type, extension):
            text = self._decode_text_content(content)
            return self._sanitize_text(text)

        data = self._decode_binary_content(content)
        if data is None:
            return self._sanitize_text(content)

        if mime_type == "application/pdf" or extension == ".pdf":
            return self._extract_pdf(data)
        if mime_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ] or extension in [".docx", ".doc"]:
            return self._extract_docx(data)
        if mime_type in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ] or extension in [".xlsx", ".xls"]:
            return self._extract_xlsx(data)
        if mime_type in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ] or extension in [".pptx", ".ppt"]:
            return self._extract_pptx(data)

        return f"[Unsupported file format: {filename}]"

    def _is_text_like(self, mime_type: str, extension: str) -> bool:
        if mime_type.startswith("text/"):
            return True
        if mime_type in [
            "application/json",
            "application/xml",
            "application/x-yaml",
            "application/yaml",
            "application/toml",
        ]:
            return True
        return extension in {
            ".json",
            ".xml",
            ".yaml",
            ".yml",
            ".toml",
            ".txt",
            ".md",
            ".markdown",
            ".rst",
            ".csv",
        }

    def _decode_text_content(self, content: str) -> str:
        try:
            decoded = base64.b64decode(content, validate=True)
            return decoded.decode("utf-8")
        except Exception:
            return content

    def _decode_binary_content(self, content: str) -> Optional[bytes]:
        try:
            return base64.b64decode(content, validate=True)
        except Exception:
            return None

    def _sanitize_text(self, text: str) -> str:
        cleaned = text.replace("\x00", "")
        return cleaned[: self.MAX_TEXT_LENGTH]

    def _extract_pdf(self, data: bytes) -> str:
        text_parts: list[str] = []
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                for i, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"--- Page {i + 1} ---\n{text}")
        except Exception as exc:
            return f"[Failed to extract PDF: {exc}]"
        return self._sanitize_text("\n\n".join(text_parts))

    def _extract_docx(self, data: bytes) -> str:
        try:
            doc = Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        except Exception as exc:
            return f"[Failed to extract Word document: {exc}]"
        return self._sanitize_text("\n\n".join(paragraphs))

    def _extract_xlsx(self, data: bytes) -> str:
        text_parts: list[str] = []
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    text_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()
        except Exception as exc:
            return f"[Failed to extract spreadsheet: {exc}]"
        return self._sanitize_text("\n\n".join(text_parts))

    def _extract_pptx(self, data: bytes) -> str:
        text_parts: list[str] = []
        try:
            prs = Presentation(io.BytesIO(data))
            for i, slide in enumerate(prs.slides):
                slide_text: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
        except Exception as exc:
            return f"[Failed to extract presentation: {exc}]"
        return self._sanitize_text("\n\n".join(text_parts))
