"""Tests for file extraction service."""

import base64
import io

import fitz
from docx import Document
import openpyxl
from pptx import Presentation
from pptx.util import Inches

from janus_gateway.services.file_extractor import FileExtractor


def _encode(data: bytes) -> str:
    return base64.b64encode(data).decode("utf-8")


def test_extract_text_file() -> None:
    extractor = FileExtractor()
    content = _encode(b"Hello text file")
    extracted = extractor.extract(content, "text/plain", "notes.txt")
    assert "Hello text file" in extracted


def test_extract_pdf() -> None:
    extractor = FileExtractor()
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF")
    pdf_bytes = doc.tobytes()
    extracted = extractor.extract(_encode(pdf_bytes), "application/pdf", "sample.pdf")
    assert "Hello PDF" in extracted


def test_extract_docx() -> None:
    extractor = FileExtractor()
    doc = Document()
    doc.add_paragraph("Hello Word")
    buffer = io.BytesIO()
    doc.save(buffer)
    extracted = extractor.extract(
        _encode(buffer.getvalue()),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "sample.docx",
    )
    assert "Hello Word" in extracted


def test_extract_xlsx() -> None:
    extractor = FileExtractor()
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet.append(["Header", "Value"])
    sheet.append(["A", 1])
    buffer = io.BytesIO()
    workbook.save(buffer)
    extracted = extractor.extract(
        _encode(buffer.getvalue()),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "sample.xlsx",
    )
    assert "Sheet1" in extracted
    assert "Header" in extracted


def test_extract_pptx() -> None:
    extractor = FileExtractor()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    textbox.text_frame.text = "Hello Slide"
    buffer = io.BytesIO()
    presentation.save(buffer)
    extracted = extractor.extract(
        _encode(buffer.getvalue()),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "sample.pptx",
    )
    assert "Hello Slide" in extracted
